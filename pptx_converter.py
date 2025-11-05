import io
import math
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

# standaard layouts
TITLE_LAYOUT = 0
TITLE_ONLY_LAYOUT = 5
BLANK_LAYOUT = 6

# instellingen
MAX_LINES_PER_SLIDE = 12
CHARS_PER_LINE = 75
MAX_BOTTOM_INCH = 6.6  # tot waar we willen schrijven op de dia


# ---------- DOCX helpers ----------
def extract_images(doc):
    """alle afbeeldingen in het docx in volgorde"""
    imgs = []
    idx = 1
    for rel in doc.part.rels.values():
        if rel.reltype == RT.IMAGE:
            blob = rel.target_part.blob
            ext = rel.target_part.partname.ext
            filename = f"image_{idx}.{ext}"
            imgs.append((filename, blob))
            idx += 1
    return imgs


def is_word_list_paragraph(para):
    """herken opsomming uit Word"""
    name = (para.style.name or "").lower()
    if "list" in name or "lijst" in name or "opsom" in name:
        return True
    ppr = para._p.pPr
    return ppr is not None and ppr.numPr is not None


def has_bold(para):
    return any(run.bold for run in para.runs)


def para_text_plain(para):
    return "".join(run.text for run in para.runs if run.text).strip()


def estimate_line_count(text: str) -> int:
    if not text:
        return 0
    return max(1, math.ceil(len(text) / CHARS_PER_LINE))


# ---------- PPTX helpers ----------
def make_bullet(paragraph):
    """
    Echte PowerPoint-bullet met wat extra ruimte (~8 mm).
    Geen qn(...) gebruiken om pptx-corruptie te voorkomen.
    """
    pPr = paragraph._p.get_or_add_pPr()

    # verwijder eventueel 'geen bullet'
    for child in list(pPr):
        if child.tag.endswith("buNone"):
            pPr.remove(child)

    # 8 mm ≈ 288000 EMU
    pPr.set("marL", "288000")     # waar de tekst begint
    pPr.set("indent", "-144000")  # bullet iets naar links

    buChar = OxmlElement("a:buChar")
    buChar.set("char", "•")
    pPr.append(buChar)


def add_textbox(slide, text, top_inch=1.0, est_lines=1):
    """voor gewone alinea's"""
    left = Inches(0.8)
    top = Inches(top_inch)
    width = Inches(8.0)
    height_inch = 0.6 + (est_lines - 1) * 0.25
    height_inch = min(height_inch, 4.0)

    shape = slide.shapes.add_textbox(left, top, width, Inches(height_inch))
    tf = shape.text_frame
    tf.text = text
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)

    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0, 0, 0)

    return height_inch


def create_title_slide(prs, title="Inhoud uit Word"):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Geconverteerd voor LessonUp"
    return slide


def create_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY_LAYOUT])
    slide.shapes.title.text = title_text
    return slide


def create_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])


def add_inline_image(slide, img_bytes, top_inch):
    """zet afbeelding op huidige dia, vaste plek"""
    left = Inches(1.0)
    top = Inches(top_inch)
    width = Inches(4.5)
    slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width)
    return 3.0  # aangenomen hoogte


# ---------- MAIN ----------
def docx_to_pptx(file_like):
    doc = Document(file_like)
    prs = Presentation()

    all_images = extract_images(doc)
    img_ptr = 0

    current_slide = create_title_slide(prs)
    current_y = 2.0
    used_lines = 0

    # state voor een lopende lijst
    current_list_tf = None
    current_list_top = 0.0
    current_list_lines = 0

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        has_image = any("graphic" in run._element.xml for run in para.runs)

        is_heading = para.style.name.startswith("Heading")
        is_bold_title = has_bold(para) and not has_image and not is_word_list_paragraph(para)
        is_list = is_word_list_paragraph(para)

        # uit een lijst lopen
        if not is_list:
            current_list_tf = None
            current_list_lines = 0

        # 1. kop / vet → nieuwe dia
        if is_heading or is_bold_title:
            current_slide = create_title_only_slide(prs, para_text_plain(para))
            current_y = 2.0
            used_lines = 0
            current_list_tf = None
            continue

        # 2. afbeelding → op deze dia als er plek is
        if has_image:
            if img_ptr < len(all_images):
                _, img_bytes = all_images[img_ptr]
                img_ptr += 1
                if current_y + 3.0 <= MAX_BOTTOM_INCH:
                    h_used = add_inline_image(current_slide, img_bytes, current_y)
                    current_y += h_used + 0.2
                else:
                    current_slide = create_blank_slide(prs)
                    h_used = add_inline_image(current_slide, img_bytes, 1.0)
                    current_y = 1.0 + h_used + 0.2
                    used_lines = 0
            continue

        # 3. lijst → één textbox, echte bullets, wrap
        if is_list and text:
            lines_needed = estimate_line_count(text)

            # past niet meer → nieuwe dia
            if (
                used_lines + lines_needed > MAX_LINES_PER_SLIDE
                or current_y + 0.6 > MAX_BOTTOM_INCH
                or (current_list_tf is None and used_lines >= MAX_LINES_PER_SLIDE)
            ):
                current_slide = create_blank_slide(prs)
                current_y = 1.0
                used_lines = 0
                current_list_tf = None
                current_list_lines = 0

            if current_list_tf is None:
                # nieuw lijstvak
                left = Inches(0.8)
                top = Inches(current_y)
                width = Inches(7.0)   # iets smaller zodat lange bullets passen
                height = Inches(4.0)
                shape = current_slide.shapes.add_textbox(left, top, width, height)
                tf = shape.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.margin_left = Inches(0.1)
                tf.margin_right = Inches(0.1)

                p = tf.paragraphs[0]
                p.text = text
                make_bullet(p)
                for r in p.runs:
                    r.font.name = "Arial"
                    r.font.size = Pt(16)

                current_list_tf = tf
                current_list_top = current_y
                current_list_lines = lines_needed
            else:
                p = current_list_tf.add_paragraph()
                p.text = text
                make_bullet(p)
                for r in p.runs:
                    r.font.name = "Arial"
                    r.font.size = Pt(16)
                current_list_lines += lines_needed

            used_lines += lines_needed
            # y onder het blok zetten
            current_y = current_list_top + 0.35 * current_list_lines + 0.3
            continue

        # 4. gewone tekst
        if text:
            lines_needed = estimate_line_count(text)
            if (
                used_lines + lines_needed > MAX_LINES_PER_SLIDE
                or current_y + 0.6 > MAX_BOTTOM_INCH
            ):
                current_slide = create_blank_slide(prs)
                current_y = 1.0
                used_lines = 0

            h_used = add_textbox(current_slide, text, top_inch=current_y, est_lines=lines_needed)
            current_y += h_used + 0.15
            used_lines += lines_needed

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out

