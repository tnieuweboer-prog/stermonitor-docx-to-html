import streamlit as st
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

st.set_page_config(page_title="DOCX → PowerPoint")
st.title("DOCX → PowerPoint (LessonUp)")

uploaded = st.file_uploader("Upload Word-bestand (.docx)", type=["docx"])


def extract_images(doc: Document):
    images = []
    idx = 1
    for rel in doc.part.rels.values():
        if rel.reltype == RT.IMAGE:
            blob = rel.target_part.blob
            ext = rel.target_part.partname.ext
            filename = f"image_{idx}.{ext}"
            images.append((filename, blob))
            idx += 1
    return images


def is_word_list_paragraph(p):
    name = (p.style.name or "").lower()
    if "list" in name or "lijst" in name or "opsom" in name:
        return True
    ppr = p._p.pPr
    return ppr is not None and ppr.numPr is not None


def get_body(slide, prs):
    """geef (text_frame, slide) terug, desnoods door nieuwe slide te maken"""
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            return shape.text_frame, slide
    # maak nieuwe slide met body
    new_slide = prs.slides.add_slide(prs.slide_layouts[1])
    return new_slide.shapes.placeholders[1].text_frame, new_slide


def apply_text_style(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0, 0, 0)


def docx_to_pptx(file):
    doc = Document(file)
    prs = Presentation()

    all_images = extract_images(doc)
    img_i = 0

    # eerste slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Inhoud uit Word"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Geconverteerd voor LessonUp"
    current_slide = slide

    for p in doc.paragraphs:
        text = (p.text or "").strip()
        has_image = any("graphic" in r._element.xml for r in p.runs)

        # nieuwe dia bij heading
        if p.style.name.startswith("Heading"):
            current_slide = prs.slides.add_slide(prs.slide_layouts[1])
            current_slide.shapes.title.text = text
            # body leeg
            body_tf, current_slide = get_body(current_slide, prs)
            body_tf.text = ""
            apply_text_style(body_tf)
            continue

        # afbeelding → aparte slide
        if has_image:
            img_slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
            img_slide.shapes.title.text = "Afbeelding"
            if img_i < len(all_images):
                _, img_bytes = all_images[img_i]
                img_i += 1
                img_stream = io.BytesIO(img_bytes)
                img_slide.shapes.add_picture(img_stream, Inches(1), Inches(1.2), width=Inches(6))
            current_slide = img_slide
            continue

        # opsomming
        if is_word_list_paragraph(p):
            body_tf, current_slide = get_body(current_slide, prs)
            para = body_tf.add_paragraph()
            para.text = text
            para.level = 0
            apply_text_style(body_tf)
            continue

        # gewone tekst
        if text:
            body_tf, current_slide = get_body(current_slide, prs)
            if body_tf.text == "":
                body_tf.text = text
            else:
                para = body_tf.add_paragraph()
                para.text = text
                para.level = 0
            apply_text_style(body_tf)

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


if uploaded:
    pptx_bytes = docx_to_pptx(uploaded)
    st.download_button(
        "Download PowerPoint (.pptx)",
        data=pptx_bytes,
        file_name="lessonup_import.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
else:
    st.info("Upload een .docx om een PowerPoint te maken.")

