import io
import os
import re
import json
import requests
from copy import deepcopy

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


# =========================
# CONFIG
# =========================
BASE_TEMPLATE_NAME = "basis layout.pptx"  # in ./templates/
LOCAL_LOGO_PATH = os.path.join(os.path.dirname(__file__), "assets", "logo.png")

# LLM config via env
LLM_PROVIDER = os.getenv("LLM_PROVIDER", "OLLAMA").upper()           # OLLAMA | OPENAI_COMPAT
LLM_MODEL    = os.getenv("LLM_MODEL", "mistral")                      # bv. 'mistral', 'qwen2.5:7b-instruct', 'llama3.1'
LLM_BASE_URL = os.getenv("LLM_BASE_URL", "http://localhost:11434")    # Ollama default; voor OPENAI_COMPAT bv. http://localhost:1234/v1
LLM_API_KEY  = os.getenv("LLM_API_KEY")                               # alleen voor OPENAI_COMPAT indien nodig


# =========================
# 0. LLM Client (zonder OpenAI SDK)
# =========================
class LLMError(RuntimeError):
    pass


class LLMClient:
    """
    Minimale client met 2 providers:
    - OLLAMA (chat API): POST /api/chat  (support 'format': 'json' -> valide JSON)
    - OPENAI_COMPAT: POST /chat/completions  (LM Studio / vLLM / andere compatibele servers)
    Geeft JSON-string terug.
    """

    def __init__(self, provider: str, model: str, base_url: str, api_key: str | None):
        self.provider = provider.upper()
        self.model = model
        self.base_url = base_url.rstrip("/")
        self.api_key = api_key

    def chat_json(self, user_prompt: str) -> str:
        if self.provider == "OLLAMA":
            return self._chat_ollama(user_prompt)
        elif self.provider == "OPENAI_COMPAT":
            return self._chat_openai_compat(user_prompt)
        else:
            raise LLMError(f"Onbekende LLM_PROVIDER: {self.provider}")

    def _chat_ollama(self, user_prompt: str) -> str:
        """
        Ollama chat API:
        POST {base}/api/chat
        body: {model, messages, stream=false, options?, format='json'}
        return: response['message']['content']
        """
        url = f"{self.base_url}/api/chat"
        payload = {
            "model": self.model,
            "messages": [{"role": "user", "content": user_prompt}],
            "stream": False,
            "format": "json",  # dwing JSON af bij veel modellen
        }
        try:
            r = requests.post(url, json=payload, timeout=120)
            r.raise_for_status()
            data = r.json()
            content = (data.get("message") or {}).get("content")
            if not content:
                raise LLMError("Lege content van Ollama.")
            return content
        except requests.RequestException as e:
            raise LLMError(f"Ollama call faalde: {e}") from e
        except ValueError:
            raise LLMError("Ollama gaf geen JSON terug.")

    def _chat_openai_compat(self, user_prompt: str) -> str:
        """
        OpenAI-compatible /chat/completions.
        Probeer response_format=json_object indien ondersteund.
        """
        url = f"{self.base_url}/chat/completions"
        headers = {"Content-Type": "application/json"}
        if self.api_key:
            headers["Authorization"] = f"Bearer {self.api_key}"

        payload = {
            "model": self.model,
            "messages": [{"role": "user", "content": user_prompt}],
            # veel compat-servers ondersteunen dit, zo niet: content bevat JSON als tekst
            "response_format": {"type": "json_object"},
            "temperature": 0.2,
        }
        try:
            r = requests.post(url, json=payload, headers=headers, timeout=120)
            r.raise_for_status()
            data = r.json()
            content = data["choices"][0]["message"]["content"]
            if not content:
                raise LLMError("Lege content van OpenAI-compat endpoint.")
            return content
        except requests.RequestException as e:
            raise LLMError(f"OpenAI-compat call faalde: {e}") from e
        except (KeyError, ValueError):
            raise LLMError("OpenAI-compat gaf onverwachte payload.")


def force_json_or_raise(text: str) -> dict:
    """
    Probeer JSON te parsen; als er ruis omheen staat, strip dan tot { ... }.
    """
    try:
        return json.loads(text)
    except Exception:
        # probeer JSON uit vrije tekst te vissen
        start = text.find("{")
        end = text.rfind("}")
        if start != -1 and end != -1 and end > start:
            snippet = text[start : end + 1]
            return json.loads(snippet)
        raise


# =========================
# 1. DOCX → blokken (kop + tekst)
# =========================
def docx_to_blocks(doc: Document) -> list[dict]:
    """
    Structuur uit Word:
    elke heading / vet / ALL CAPS = nieuwe dia
    onderliggende tekst = body
    return: [{"title": "...", "body": "..."}, ...]
    """
    blocks = []
    current_title = None
    current_body: list[str] = []

    for para in doc.paragraphs:
        txt = "".join(r.text for r in para.runs if r.text).strip()
        if not txt:
            continue

        is_heading = (
            (para.style and para.style.name and para.style.name.lower().startswith("heading"))
            or any(r.bold for r in para.runs)
            or (len(txt) <= 50 and txt.upper() == txt)  # korte regel in CAPS
        )

        if is_heading:
            if current_title or current_body:
                blocks.append({"title": current_title, "body": "\n".join(current_body).strip()})
            current_title = txt
            current_body = []
        else:
            current_body.append(txt)

    if current_title or current_body:
        blocks.append({"title": current_title, "body": "\n".join(current_body).strip()})

    if not blocks:
        blocks = [{"title": "Lesonderdeel", "body": "(Geen duidelijke structuur gevonden in dit document.)"}]

    return blocks


# =========================
# 2. LLM (zonder OpenAI SDK): alle blokken → slides
# =========================
def llm_make_all_slides_from_blocks(blocks: list[dict]) -> list[dict]:
    """
    Stuurt ALLE blokken in één prompt naar het gekozen model (Ollama / OpenAI-compat).
    Return: [{"title":"...","text":["...","..."],"check":"..."}...]
    """
    client = LLMClient(LLM_PROVIDER, LLM_MODEL, LLM_BASE_URL, LLM_API_KEY)

    parts = []
    for i, b in enumerate(blocks, start=1):
        parts.append(f"### Onderdeel {i}\nKop: {b.get('title') or ''}\nTekst:\n{b.get('body') or ''}\n")
    joined = "\n\n".join(parts)

    prompt = f"""
Je krijgt hieronder meerdere onderdelen uit een les over installatietechniek.
Maak hier dia's van voor een VMBO-les (basis/kader/GL).

Voor elk onderdeel:
- bedenk 1 korte, begrijpelijke titel (max 8 woorden)
- schrijf 2 of 3 korte, vertellende zinnen in de je-vorm
- schrijf 1 controlevraag die past bij de uitleg
- herhaal de titel NIET in de tekst
- gebruik eenvoudige woorden

Geef ALLEEN geldig JSON in dit formaat:

{{
  "slides": [
    {{
      "title": "…",
      "text": ["…", "…", "…"],
      "check": "…"
    }}
  ]
}}

Hier zijn de onderdelen uit het Word-document:

{joined}
"""
    try:
        raw = client.chat_json(prompt)
        data = force_json_or_raise(raw)
    except Exception as e:
        raise LLMError(f"LLM gaf geen geldig JSON: {e}") from e

    slides = data.get("slides")
    if not slides or not isinstance(slides, list):
        raise LLMError("LLM antwoord bevat geen 'slides' lijst.")
    # minimale normalisatie
    norm = []
    for s in slides:
        title = (s.get("title") or "Lesonderdeel").strip()
        text = s.get("text") or []
        if isinstance(text, str):
            text = [text]
        text = [t.strip() for t in text if t and t.strip()]
        check = (s.get("check") or "").strip()
        norm.append({"title": title, "text": text, "check": check})
    return norm


# =========================
# 3. Fallback: zonder LLM → heuristisch
# =========================
def fallback_slides_from_blocks(blocks: list[dict]) -> list[dict]:
    slides = []
    for b in blocks:
        title_raw = (b.get("title") or "Lesonderdeel").strip()
        body_raw = (b.get("body") or "").strip()
        title = title_raw.capitalize()

        sentences = re.split(r"[.!?]\s+", body_raw)
        sentences = [s.strip(" .!?") for s in sentences if s.strip()]
        text_lines = []
        for s in sentences[:3]:
            s = s.replace("Men ", "Je ").replace("men ", "je ")
            if not s.lower().startswith(("je ", "dit ", "zo ", "dan ", "als ")):
                s = "Je " + s[0].lower() + s[1:]
            text_lines.append(s)

        if not text_lines:
            text_lines = [
                "Je leert hier hoe je dit onderdeel goed uitvoert.",
                "Zo kan water en lucht goed weg.",
                "Dan krijg je geen stank.",
            ]

        lower = body_raw.lower()
        if any(k in lower for k in ("niet", "mag", "nooit")):
            check = "Waarom mag je dit niet zo doen?"
        elif any(k in lower for k in ("leiding", "afvoer")):
            check = "Wat gebeurt er als je dit verkeerd aansluit?"
        else:
            check = "Kun je uitleggen waarom je dit zo doet?"

        slides.append({"title": title, "text": text_lines, "check": check})
    return slides


# =========================
# 4. PPTX helpers
# =========================
def get_logo_bytes():
    if os.path.exists(LOCAL_LOGO_PATH):
        with open(LOCAL_LOGO_PATH, "rb") as f:
            return f.read()
    return None


def add_logo(slide, logo_bytes):
    if not logo_bytes:
        return
    slide.shapes.add_picture(io.BytesIO(logo_bytes), Inches(7.5), Inches(0.2), width=Inches(1.5))


def get_positions_from_first_slide(slide):
    text_shapes = [s for s in slide.shapes if hasattr(s, "text") and s.text and s.text.strip()]
    if len(text_shapes) >= 2:
        t, b = text_shapes[0], text_shapes[1]
        return {
            "title": {"left": t.left, "top": t.top, "width": t.width, "height": t.height},
            "body": {"left": b.left, "top": b.top, "width": b.width, "height": b.height},
        }
    return {
        "title": {"left": Inches(0.8), "top": Inches(0.8), "width": Inches(9), "height": Inches(1)},
        "body": {"left": Inches(0.8), "top": Inches(3), "width": Inches(10), "height": Inches(3)},
    }


def duplicate_slide_clean(prs: Presentation, slide_index: int):
    src = prs.slides[slide_index]
    dest = prs.slides.add_slide(prs.slide_layouts[0])
    for shp in src.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        new_el = deepcopy(shp.element)
        dest.shapes._spTree.insert_element_before(new_el, "p:extLst")
    for shp in dest.shapes:
        if hasattr(shp, "text_frame"):
            shp.text_frame.clear()
    return dest


def place_title(slide, text: str, pos: dict):
    box = slide.shapes.add_textbox(pos["left"], pos["top"], pos["width"], pos["height"])
    tf = box.text_frame
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(28)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0, 0, 0)


def place_text_and_question(slide, lines: list[str], check: str, pos: dict):
    box = slide.shapes.add_textbox(pos["left"], pos["top"], pos["width"], pos["height"])
    tf = box.text_frame
    tf.word_wrap = True

    first = True
    for line in lines:
        if not line:
            continue
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        p.text = line
        first = False

    tf.add_paragraph().text = ""

    if check:
        p = tf.add_paragraph()
        p.text = check
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(16)

    for p in tf.paragraphs:
        for r in p.runs:
            if not r.font.size:
                r.font.name = "Arial"
                r.font.size = Pt(16)
                r.font.color.rgb = RGBColor(0, 0, 0)


# =========================
# 5. MAIN: DOCX → PPTX
# =========================
def docx_to_pptx_hybrid(file_like):
    # 1) template
    base_dir = os.path.dirname(__file__)
    template_path = os.path.join(base_dir, "templates", BASE_TEMPLATE_NAME)
    prs = Presentation(template_path) if os.path.exists(template_path) else Presentation()

    # 2) input
    doc = Document(file_like)
    blocks = docx_to_blocks(doc)

    # 3) LLM of fallback
    try:
        slides_data = llm_make_all_slides_from_blocks(blocks)
    except Exception:
        slides_data = fallback_slides_from_blocks(blocks)

    # 4) logo + eerste dia
    logo_bytes = get_logo_bytes()
    if not prs.slides:
        prs.slides.add_slide(prs.slide_layouts[0])
    first_slide = prs.slides[0]
    positions = get_positions_from_first_slide(first_slide)

    # eerste dia leeg
    for shp in first_slide.shapes:
        if hasattr(shp, "text_frame"):
            shp.text_frame.clear()
    if logo_bytes:
        add_logo(first_slide, logo_bytes)

    # 5) vul dia's
    first = slides_data[0]
    place_title(first_slide, first["title"], positions["title"])
    place_text_and_question(first_slide, first.get("text", []), first.get("check", ""), positions["body"])

    for sd in slides_data[1:]:
        slide = duplicate_slide_clean(prs, 0)
        if logo_bytes:
            add_logo(slide, logo_bytes)
        place_title(slide, sd["title"], positions["title"])
        place_text_and_question(slide, sd.get("text", []), sd.get("check", ""), positions["body"])

    # 6) output
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out

