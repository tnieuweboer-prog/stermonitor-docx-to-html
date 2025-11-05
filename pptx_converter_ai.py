import io
import os
import json
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TITLE_LAYOUT = 0
TITLE_ONLY_LAYOUT = 5


# ---------- hulpfuncties ----------
def extract_images(doc):
    imgs = []
    for rel in doc.part.rels.values():
        if rel.reltype == RT.IMAGE:
            imgs.append(rel.target_part.blob)
    return imgs


def local_fallback_summarize(text: str) -> dict:
    """
    Als OpenAI niet lukt: maak zelf een dia.
    - eerste zin = titel (kort)
    - max 3 bullets uit de rest
    """
    text = text.strip()
    if not text:
        return {"title": "Lesonderdeel", "bullets": ["Inhoud kon niet worden samengevat."], "image_hint": ""}

    # titel = eerste 6 woorden
    words = text.split()
    title = " ".join(words[:6])
    if len(words) > 6:
        title += "..."

    # simpele bullet-split
    parts = [p.strip() for p in text.replace("•", "\n").split("\n") if p.strip()]
    if len(parts) == 1:
        # probeer op punt te delen
        parts = [p.strip() for p in text.split(".") if p.strip()]

    bullets = []
    for p in parts:
        if len(bullets) >= 3:
            break
        bullets.append(p[:90] + ("..." if len(p) > 90 else ""))

    if not bullets:
        bullets = [text[:90] + ("..." if len(text) > 90 else "")]

    return {
        "title": title,
        "bullets": bullets,
        "image_hint": "",
    }


def call_ai_for_slide(text: str) -> dict:
    """
    Probeer OpenAI. Als er geen key is of het gaat mis → lokale fallback.
    """
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        return local_fallback_summarize(text)

    try:
        from openai import OpenAI
        client = OpenAI(api_key=api_key)

        prompt = f"""
Je bent docent elektrotechniek en maakt precies 1 PowerPoint-dia.
Doelgroep: mbo, korte zinnen.
Maak van de tekst hieronder 1 dia met:
- title (kort, max 8 woorden)
- bullets (max 4, heel kort)
Geef ALLEEN JSON terug met: title, bullets, image_hint.

Tekst:
{text}
"""
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
        )
        data = json.loads(resp.choices[0].message.content)
        # sanity
        if not data.get("title"):
            data["title"] = text[:40]
        if not data.get("bullets"):
            # als model rare output gaf → fallback bullets
            fb = local_fallback_summarize(text)
            data["bullets"] = fb["bullets"]
        return data

    except Exception:
        # bv. modelnaam fout, key fout, netwerk fout
        return local_fallback_summarize(text)


def add_ai_slide(prs: Presentation, ai_obj: dict, image_bytes: bytes | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY_LAYOUT])

    title = ai_obj.get("title") or "Lesonderdeel"
    bullets = ai_obj.get("bullets") or ["Inhoud uit tekst."]

    slide.shapes.title.text = title

    # tekstvak links
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(5.5) if image_bytes else Inches(7.5)
    height = Inches(4.0)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.text = ""

    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0, 0, 0)

    # afbeelding rechts
    if image_bytes:
        img_left = Inches(6.0)
        img_top = Inches(2.0)
        slide.shapes.add_picture(io.BytesIO(image_bytes), img_left, img_top, width=Inches(2.5))


def docx_to_pptx_ai(file_like):
    doc = Document(file_like)
    prs = Presentation()

    all_imgs = extract_images(doc)
    img_idx = 0

    # openingsdia
    first = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    first.shapes.title.text = "Les gegenereerd uit Word"
    if len(first.placeholders) > 1:
        first.placeholders[1].text = "AI / samenvatting per alinea"

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        has_image = any("graphic" in r._element.xml for r in para.runs)

        if not text and not has_image:
            continue

        ai_obj = call_ai_for_slide(text if text else "Afbeelding bij uitleg")

        img_bytes = None
        if has_image and img_idx < len(all_imgs):
            img_bytes = all_imgs[img_idx]
            img_idx += 1

        add_ai_slide(prs, ai_obj, image_bytes=img_bytes)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out
